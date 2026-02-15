#!/usr/bin/env python3
"""Generate a polished PowerPoint pitch deck for Sustainability Signals."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG       = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
CARD_BG       = RGBColor(0x16, 0x21, 0x3A)   # Slightly lighter navy
ACCENT        = RGBColor(0x34, 0xD3, 0x99)   # Emerald green accent
ACCENT_DIM    = RGBColor(0x22, 0xA0, 0x73)   # Darker green
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xA0, 0xAE, 0xC0)
MID_GRAY      = RGBColor(0x6B, 0x7B, 0x93)
GOLD          = RGBColor(0xFA, 0xCC, 0x15)
ORANGE        = RGBColor(0xF9, 0x73, 0x16)
BLUE_ACCENT   = RGBColor(0x38, 0xBD, 0xF8)
RED_ACCENT    = RGBColor(0xF4, 0x3F, 0x5E)

SLIDE_WIDTH   = Inches(13.333)
SLIDE_HEIGHT  = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ── Helpers ─────────────────────────────────────────────────────────────────

def _solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color, line_color=None, line_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w or Pt(1)
    else:
        shape.line.fill.background()
    # Make corners subtle
    shape.adjustments[0] = 0.04
    return shape


def _add_text(slide, left, top, width, height, text, font_size=18,
              color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def _add_bullet_frame(slide, left, top, width, height, items,
                      font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT, spacing=Pt(8)):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txbox


def _pill(slide, left, top, text, fill=ACCENT, text_color=DARK_BG, width=None, font_size=11):
    w = width or Inches(1.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.adjustments[0] = 0.5  # full round
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return shape


def _section_tag(slide, left, top, label):
    _add_text(slide, left, top, Inches(4), Inches(0.35), label.upper(),
              font_size=11, color=ACCENT, bold=True)


def _slide_number(slide, num):
    _add_text(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.3),
              str(num), font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


# ── SLIDE 1 – TITLE ────────────────────────────────────────────────────────
def slide_title():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _solid_bg(slide, DARK_BG)

    # Decorative accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1), Inches(1.8), Inches(0.06), Inches(1.5))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    _add_text(slide, Inches(1.4), Inches(1.8), Inches(10), Inches(0.9),
              "SustainabilitySignals", font_size=52, color=WHITE, bold=True)

    _add_text(slide, Inches(1.4), Inches(2.7), Inches(10), Inches(0.6),
              "Transparent ESG Ratings from Disclosure Evidence", font_size=24, color=ACCENT)

    _add_text(slide, Inches(1.4), Inches(3.7), Inches(8), Inches(0.8),
              "AI-powered platform that reads 950+ European sustainability reports,\n"
              "scores disclosure quality, and extracts structured ESG data — all evidence-grounded.",
              font_size=16, color=LIGHT_GRAY)

    # Bottom bar
    _add_text(slide, Inches(1.4), Inches(5.8), Inches(5), Inches(0.4),
              "Nikhil Reddy Gogu  ·  nikhil.chat  ·  Feb 2026",
              font_size=13, color=MID_GRAY)

    _pill(slide, Inches(9.5), Inches(5.82), "Student Project", fill=CARD_BG, text_color=ACCENT, width=Inches(1.8))

    _slide_number(slide, 1)


# ── SLIDE 2 – THE PROBLEM ──────────────────────────────────────────────────
def slide_problem():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _section_tag(slide, Inches(1), Inches(0.6), "The Problem")
    _add_text(slide, Inches(1), Inches(1.1), Inches(11), Inches(0.8),
              "ESG ratings are a black box", font_size=40, bold=True)

    problems = [
        ("Opaque Methodologies",
         "Major ESG rating agencies (MSCI, Sustainalytics, ISS) use proprietary, non-reproducible scoring. Investors can't verify how a rating was derived.",
         RED_ACCENT),
        ("Low Correlation",
         "Academic research shows correlations as low as 0.38-0.71 between major ESG rating providers — the same company gets wildly different scores.",
         ORANGE),
        ("No Evidence Trail",
         "Current ratings provide a number, not evidence. When a company scores 'AA', stakeholders can't trace what disclosure drove that result.",
         GOLD),
    ]

    for i, (title, desc, accent) in enumerate(problems):
        left = Inches(1) + Inches(i * 3.8)
        card = _add_rect(slide, left, Inches(2.5), Inches(3.5), Inches(3.8), CARD_BG, accent, Pt(1.5))

        # Accent dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), Inches(2.85), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = accent; dot.line.fill.background()

        _add_text(slide, left + Inches(0.3), Inches(3.15), Inches(2.9), Inches(0.5),
                  title, font_size=20, bold=True, color=WHITE)
        _add_text(slide, left + Inches(0.3), Inches(3.7), Inches(2.9), Inches(2.2),
                  desc, font_size=14, color=LIGHT_GRAY)

    _slide_number(slide, 2)


# ── SLIDE 3 – SOLUTION ─────────────────────────────────────────────────────
def slide_solution():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _section_tag(slide, Inches(1), Inches(0.6), "Our Solution")
    _add_text(slide, Inches(1), Inches(1.1), Inches(11), Inches(0.8),
              "Evidence-grounded ESG intelligence", font_size=40, bold=True)

    _add_text(slide, Inches(1), Inches(2.1), Inches(10), Inches(0.6),
              "SustainabilitySignals reads the actual report, extracts evidence, and shows its work.",
              font_size=17, color=LIGHT_GRAY)

    pillars = [
        ("📄", "Read", "AI converts 950+ sustainability PDFs to structured text. Every page, every table."),
        ("🔍", "Score", "Regex + NLP engine detects 50+ disclosure signals across 4 quality dimensions."),
        ("🧠", "Extract", "FinBERT-ESG classifies content; LangExtract pulls structured entities (emissions, targets, policies)."),
        ("💬", "Chat", "Ask questions about any report with AI grounded in the actual document text — not hallucinations."),
    ]

    for i, (icon, label, desc) in enumerate(pillars):
        left = Inches(1) + Inches(i * 3)
        _add_rect(slide, left, Inches(3.1), Inches(2.7), Inches(3.2), CARD_BG)

        _add_text(slide, left + Inches(0.3), Inches(3.3), Inches(0.6), Inches(0.6),
                  icon, font_size=28)
        _add_text(slide, left + Inches(0.3), Inches(3.95), Inches(2.1), Inches(0.4),
                  label, font_size=22, bold=True, color=ACCENT)
        _add_text(slide, left + Inches(0.3), Inches(4.45), Inches(2.1), Inches(1.6),
                  desc, font_size=13, color=LIGHT_GRAY)

    _slide_number(slide, 3)


# ── SLIDE 4 – DISCLOSURE QUALITY ENGINE ────────────────────────────────────
def slide_dq_engine():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _section_tag(slide, Inches(1), Inches(0.6), "Core Engine")
    _add_text(slide, Inches(1), Inches(1.1), Inches(11), Inches(0.8),
              "Disclosure Quality Scoring", font_size=40, bold=True)

    # Formula card
    _add_rect(slide, Inches(1), Inches(2.3), Inches(11.3), Inches(1.2), CARD_BG, ACCENT, Pt(1))
    _add_text(slide, Inches(1.4), Inches(2.45), Inches(10), Inches(0.4),
              "DQ Score =  0.35 × Completeness  +  0.25 × Consistency  +  0.20 × Assurance  +  0.20 × Transparency",
              font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1.4), Inches(2.95), Inches(10), Inches(0.35),
              "Bands:  High ≥ 75   ·   Medium 50-74   ·   Low < 50   |   All evidence-grounded with source page references",
              font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    dims = [
        ("Completeness", "35%", "Breadth across frameworks (GRI, ESRS, TCFD, SASB), materiality, governance, climate, social, taxonomy.",
         ACCENT),
        ("Consistency", "25%", "Methodology, base year, GHG protocol, comparative data, quantitative density, data quality signals.",
         BLUE_ACCENT),
        ("Assurance", "20%", "Limited/reasonable assurance levels, standards (ISAE, AA1000), named providers, scope breadth.",
         GOLD),
        ("Transparency", "20%", "Forward-looking statements, transition plans, stakeholder engagement, financial connectivity, ESRS datapoints.",
         ORANGE),
    ]

    for i, (name, weight, desc, color) in enumerate(dims):
        left = Inches(1) + Inches(i * 2.9)
        _add_rect(slide, left, Inches(4.0), Inches(2.65), Inches(3.0), CARD_BG, color, Pt(1.5))

        _pill(slide, left + Inches(0.2), Inches(4.25), weight, fill=color, text_color=DARK_BG, width=Inches(0.7), font_size=13)
        _add_text(slide, left + Inches(1.05), Inches(4.22), Inches(1.5), Inches(0.35),
                  name, font_size=17, bold=True, color=WHITE)
        _add_text(slide, left + Inches(0.2), Inches(4.8), Inches(2.25), Inches(2.0),
                  desc, font_size=12, color=LIGHT_GRAY)

    _slide_number(slide, 4)


# ── SLIDE 5 – AI & ML PIPELINE ─────────────────────────────────────────────
def slide_ai_pipeline():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _section_tag(slide, Inches(1), Inches(0.6), "AI & ML")
    _add_text(slide, Inches(1), Inches(1.1), Inches(11), Inches(0.8),
              "Multi-model intelligence pipeline", font_size=40, bold=True)

    steps = [
        ("PDF → Markdown", "Workers AI\ntoMarkdown()", "Converts report PDFs into\nclean, structured text"),
        ("ESG Classification", "FinBERT-ESG\n9 Categories", "Routes content chunks\nby ESG topic"),
        ("Entity Extraction", "LangExtract\n+ LLM", "Pulls emissions, targets,\npolicies, metrics"),
        ("Semantic Search", "Vectorize\n+ Embeddings", "Enables relevance-ranked\nretrieval across reports"),
        ("Grounded Chat", "Mistral / Llama\nContext Window", "Q&A over report text,\nno hallucinations"),
    ]

    for i, (title, model, desc) in enumerate(steps):
        left = Inches(0.6) + Inches(i * 2.5)
        _add_rect(slide, left, Inches(2.5), Inches(2.25), Inches(3.8), CARD_BG)

        # Step number circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.8), Inches(2.75), Inches(0.5), Inches(0.5))
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT; circ.line.fill.background()
        tf = circ.text_frame; p = tf.paragraphs[0]
        p.text = str(i + 1); p.font.size = Pt(18); p.font.bold = True
        p.font.color.rgb = DARK_BG; p.alignment = PP_ALIGN.CENTER

        _add_text(slide, left + Inches(0.2), Inches(3.45), Inches(1.85), Inches(0.5),
                  title, font_size=16, bold=True, color=WHITE)
        _add_text(slide, left + Inches(0.2), Inches(3.95), Inches(1.85), Inches(0.6),
                  model, font_size=12, color=ACCENT)
        _add_text(slide, left + Inches(0.2), Inches(4.7), Inches(1.85), Inches(1.2),
                  desc, font_size=12, color=LIGHT_GRAY)

        # Arrow between steps
        if i < len(steps) - 1:
            _add_text(slide, left + Inches(2.25), Inches(3.7), Inches(0.3), Inches(0.5),
                      "→", font_size=24, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    _slide_number(slide, 5)


# ── SLIDE 6 – COVERAGE UNIVERSE ────────────────────────────────────────────
def slide_coverage():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _section_tag(slide, Inches(1), Inches(0.6), "Coverage")
    _add_text(slide, Inches(1), Inches(1.1), Inches(11), Inches(0.8),
              "953 European sustainability reports", font_size=40, bold=True)

    _add_text(slide, Inches(1), Inches(2.1), Inches(10), Inches(0.5),
              "CSRD-aligned disclosures from 2024 annual & sustainability reports, covering 32 countries and 11 GICS sectors.",
              font_size=16, color=LIGHT_GRAY)

    # Big stats
    stats = [
        ("953", "Reports"),
        ("32", "Countries"),
        ("11", "GICS Sectors"),
        ("50+", "DQ Features"),
    ]
    for i, (num, label) in enumerate(stats):
        left = Inches(1) + Inches(i * 3)
        _add_rect(slide, left, Inches(3.0), Inches(2.65), Inches(1.6), CARD_BG, ACCENT, Pt(1))
        _add_text(slide, left, Inches(3.15), Inches(2.65), Inches(0.7),
                  num, font_size=42, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
        _add_text(slide, left, Inches(3.85), Inches(2.65), Inches(0.4),
                  label, font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Top sectors table
    _add_text(slide, Inches(1), Inches(5.0), Inches(5), Inches(0.4),
              "Top Sectors by Coverage", font_size=18, bold=True)

    sectors = [
        ("Industrials", "233"), ("Financials", "162"), ("Consumer Discretionary", "88"),
        ("Information Technology", "85"), ("Materials", "81"), ("Consumer Staples", "73"),
    ]
    for i, (sector, count) in enumerate(sectors):
        y = Inches(5.5) + Inches(i * 0.3)
        _add_text(slide, Inches(1.2), y, Inches(3), Inches(0.3),
                  f"  {sector}", font_size=12, color=LIGHT_GRAY)
        _add_text(slide, Inches(4), y, Inches(0.8), Inches(0.3),
                  count, font_size=12, color=ACCENT, bold=True, alignment=PP_ALIGN.RIGHT)

    # Countries
    _add_text(slide, Inches(6.5), Inches(5.0), Inches(6), Inches(0.4),
              "Pan-European Coverage", font_size=18, bold=True)
    _add_text(slide, Inches(6.5), Inches(5.5), Inches(6), Inches(2.0),
              "Germany · France · Netherlands · Italy · Spain · Belgium · Sweden\n"
              "Norway · Finland · Denmark · Austria · UK · Switzerland · Ireland\n"
              "Luxembourg · Poland · Portugal · Greece · Czech Republic · Romania\n"
              "and 12 more EU/EEA countries",
              font_size=13, color=LIGHT_GRAY)

    _slide_number(slide, 6)


# ── SLIDE 7 – PLATFORM FEATURES ────────────────────────────────────────────
def slide_features():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _section_tag(slide, Inches(1), Inches(0.6), "Platform")
    _add_text(slide, Inches(1), Inches(1.1), Inches(11), Inches(0.8),
              "Full-stack ESG intelligence platform", font_size=40, bold=True)

    features = [
        ("🔎", "Report Library", "Browse, filter, and search 953 sustainability reports by company, sector, country, or year.", "LIVE", ACCENT),
        ("📊", "DQ Scoring", "Transparent Disclosure Quality scores with subscores, evidence highlights, and source page references.", "LIVE", ACCENT),
        ("🤖", "AI Chat", "Ask questions about any report. The AI is grounded in actual document text — with page citations.", "LIVE", ACCENT),
        ("📄", "In-App PDF", "Read sustainability reports directly in the browser with range-served PDFs for instant rendering.", "LIVE", ACCENT),
        ("🧬", "Entity Extraction", "FinBERT-ESG + LangExtract pulls structured ESG entities: emissions, targets, policies, metrics.", "LIVE", ACCENT),
        ("📈", "Company Profiles", "Deep-dive ESG profiles with score breakdowns, risk factors, peer comparison, and time-series.", "PLANNED", MID_GRAY),
    ]

    for i, (icon, title, desc, status, status_color) in enumerate(features):
        col = i % 3
        row = i // 3
        left = Inches(1) + Inches(col * 3.9)
        top  = Inches(2.3) + Inches(row * 2.5)
        _add_rect(slide, left, top, Inches(3.6), Inches(2.2), CARD_BG)

        _add_text(slide, left + Inches(0.2), top + Inches(0.2), Inches(0.5), Inches(0.5),
                  icon, font_size=22)
        _pill(slide, left + Inches(2.4), top + Inches(0.25), status, fill=status_color,
              text_color=DARK_BG if status == "LIVE" else WHITE, width=Inches(1.0), font_size=9)
        _add_text(slide, left + Inches(0.2), top + Inches(0.7), Inches(3.2), Inches(0.4),
                  title, font_size=17, bold=True, color=WHITE)
        _add_text(slide, left + Inches(0.2), top + Inches(1.15), Inches(3.2), Inches(1.0),
                  desc, font_size=12, color=LIGHT_GRAY)

    _slide_number(slide, 7)


# ── SLIDE 8 – TECH STACK ───────────────────────────────────────────────────
def slide_tech():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _section_tag(slide, Inches(1), Inches(0.6), "Technology")
    _add_text(slide, Inches(1), Inches(1.1), Inches(11), Inches(0.8),
              "Cloud-native, edge-first architecture", font_size=40, bold=True)

    layers = [
        ("Frontend", [
            "React 19 + TypeScript + Vite",
            "react-pdf for in-browser PDF rendering",
            "Tailwind CSS for responsive design",
            "Client-side routing via React Router",
        ], ACCENT),
        ("Backend / Edge", [
            "Cloudflare Pages Functions (serverless)",
            "Cloudflare R2 (object storage for PDFs)",
            "Cloudflare Workers AI (Mistral, Llama)",
            "Cloudflare Vectorize (semantic search)",
        ], BLUE_ACCENT),
        ("ML / AI", [
            "FinBERT-ESG-9 (ESG classification)",
            "LangExtract (structured entity extraction)",
            "ESG-BERT (issue classification + pillar)",
            "Azure Container Apps (inference hosting)",
        ], GOLD),
        ("Data Pipeline", [
            "PDF → Markdown via AI.toMarkdown()",
            "Regex v4.1 DQ feature engine",
            "Vectorize chunk embeddings (bge-small)",
            "NDJSON manifests + R2 cache layer",
        ], ORANGE),
    ]

    for i, (title, items, color) in enumerate(layers):
        left = Inches(1) + Inches(i * 3.05)
        _add_rect(slide, left, Inches(2.3), Inches(2.8), Inches(4.4), CARD_BG, color, Pt(1))
        _add_text(slide, left + Inches(0.25), Inches(2.5), Inches(2.3), Inches(0.4),
                  title, font_size=18, bold=True, color=color)

        for j, item in enumerate(items):
            y = Inches(3.1) + Inches(j * 0.85)
            _add_text(slide, left + Inches(0.25), y, Inches(2.3), Inches(0.75),
                      item, font_size=12, color=LIGHT_GRAY)

    _slide_number(slide, 8)


# ── SLIDE 9 – TRACTION ─────────────────────────────────────────────────────
def slide_traction():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _section_tag(slide, Inches(1), Inches(0.6), "Traction")
    _add_text(slide, Inches(1), Inches(1.1), Inches(11), Inches(0.8),
              "What we've built (solo, in months)", font_size=40, bold=True)

    milestones = [
        ("953 reports", "collected, extracted, and indexed from EU-listed companies"),
        ("DQ Engine v4.1", "50+ disclosure features with evidence-grounded scoring"),
        ("Entity extraction", "structured ESG entities via FinBERT + LangExtract pipeline"),
        ("AI chat", "grounded Q&A over full report text with page citations"),
        ("Semantic search", "Vectorize-powered retrieval across entire report corpus"),
        ("Live platform", "deployed end-to-end on Cloudflare edge infrastructure"),
    ]

    for i, (metric, desc) in enumerate(milestones):
        col = i % 2
        row = i // 2
        left = Inches(1) + Inches(col * 6)
        top  = Inches(2.5) + Inches(row * 1.5)
        _add_rect(slide, left, top, Inches(5.6), Inches(1.2), CARD_BG)

        # Checkmark circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), top + Inches(0.3), Inches(0.5), Inches(0.5))
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT; circ.line.fill.background()
        tf = circ.text_frame; p = tf.paragraphs[0]
        p.text = "✓"; p.font.size = Pt(18); p.font.bold = True
        p.font.color.rgb = DARK_BG; p.alignment = PP_ALIGN.CENTER

        _add_text(slide, left + Inches(1.0), top + Inches(0.2), Inches(4.3), Inches(0.4),
                  metric, font_size=18, bold=True, color=WHITE)
        _add_text(slide, left + Inches(1.0), top + Inches(0.65), Inches(4.3), Inches(0.4),
                  desc, font_size=13, color=LIGHT_GRAY)

    _slide_number(slide, 9)


# ── SLIDE 10 – ROADMAP ─────────────────────────────────────────────────────
def slide_roadmap():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _section_tag(slide, Inches(1), Inches(0.6), "Roadmap")
    _add_text(slide, Inches(1), Inches(1.1), Inches(11), Inches(0.8),
              "From Disclosure Quality to ESG Ratings", font_size=40, bold=True)

    phases = [
        ("NOW", "Foundation", [
            "953 reports scored & indexed",
            "DQ engine v4.1 live",
            "Entity extraction operational",
            "AI chat functional",
        ], ACCENT, True),
        ("Q2 2026", "Intelligence Layer", [
            "Company ESG profiles / dashboards",
            "Peer comparison & benchmarking",
            "Cross-company entity aggregation",
            "Time-series DQ tracking",
        ], BLUE_ACCENT, False),
        ("Q3 2026", "ESG Ratings", [
            "Evidence-grounded ESG ratings",
            "Risk factor identification",
            "Industry materiality weighting",
            "API access for institutions",
        ], GOLD, False),
        ("Q4 2026", "Scale", [
            "Expand beyond EU (US, APAC)",
            "Real-time CSRD monitoring",
            "Regulatory change detection",
            "Institutional data feeds",
        ], ORANGE, False),
    ]

    for i, (when, title, items, color, is_active) in enumerate(phases):
        left = Inches(0.7) + Inches(i * 3.15)
        border = color if is_active else None
        bw = Pt(2) if is_active else None
        _add_rect(slide, left, Inches(2.3), Inches(2.9), Inches(4.6), CARD_BG, border, bw)

        _pill(slide, left + Inches(0.25), Inches(2.55), when, fill=color, text_color=DARK_BG, width=Inches(1.2), font_size=11)
        _add_text(slide, left + Inches(0.25), Inches(3.05), Inches(2.4), Inches(0.4),
                  title, font_size=19, bold=True, color=WHITE)

        for j, item in enumerate(items):
            y = Inches(3.65) + Inches(j * 0.7)
            # Small dot
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), y + Inches(0.08), Inches(0.1), Inches(0.1))
            dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
            _add_text(slide, left + Inches(0.55), y, Inches(2.2), Inches(0.55),
                      item, font_size=12, color=LIGHT_GRAY)

    _slide_number(slide, 10)


# ── SLIDE 11 – TEAM ────────────────────────────────────────────────────────
def slide_team():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _section_tag(slide, Inches(1), Inches(0.6), "Team")
    _add_text(slide, Inches(1), Inches(1.1), Inches(11), Inches(0.8),
              "Built by a finance student, for the finance world", font_size=40, bold=True)

    # Main profile card
    _add_rect(slide, Inches(1), Inches(2.5), Inches(5.5), Inches(4.2), CARD_BG, ACCENT, Pt(1))

    _add_text(slide, Inches(1.5), Inches(2.8), Inches(4.5), Inches(0.5),
              "Nikhil Reddy Gogu", font_size=28, bold=True, color=WHITE)
    _add_text(slide, Inches(1.5), Inches(3.35), Inches(4.5), Inches(0.35),
              "Founder & Builder", font_size=16, color=ACCENT)

    details = [
        "🎓  Dual Degree: Master in Finance (AMS) + MBA Financial Services (Woxsen)",
        "🇧🇪  Based in Antwerpen, Belgium",
        "📊  ESG Experience: Impactree Data Technologies",
        "🌍  PRME Global Working Group (Steering Committee)",
        "📈  Bloomberg Lab Student Representative",
        "🏫  Class of '26 Representative at Antwerp Management School",
        "💻  Full-stack engineer: React, TypeScript, Python, ML, Cloudflare",
    ]
    for i, detail in enumerate(details):
        _add_text(slide, Inches(1.5), Inches(4.0) + Inches(i * 0.36), Inches(4.8), Inches(0.35),
                  detail, font_size=12, color=LIGHT_GRAY)

    # Skills / highlights
    _add_rect(slide, Inches(7), Inches(2.5), Inches(5.3), Inches(2.0), CARD_BG)
    _add_text(slide, Inches(7.4), Inches(2.7), Inches(4.5), Inches(0.4),
              "Why this matters", font_size=18, bold=True, color=WHITE)
    _add_text(slide, Inches(7.4), Inches(3.2), Inches(4.5), Inches(1.1),
              "Finance domain expertise + engineering capability = \n"
              "a platform that understands what investors actually need\n"
              "from ESG disclosures and can build it end-to-end.",
              font_size=13, color=LIGHT_GRAY)

    _add_rect(slide, Inches(7), Inches(4.8), Inches(5.3), Inches(1.9), CARD_BG)
    _add_text(slide, Inches(7.4), Inches(5.0), Inches(4.5), Inches(0.4),
              "Links", font_size=18, bold=True, color=WHITE)
    links = [
        ("🌐  Portfolio:", "nikhil.chat"),
        ("💼  LinkedIn:", "linkedin.com/in/nikhilgogu"),
        ("📧  Email:", "nikhilreddy.gogu@student.ams.ac.be"),
    ]
    for i, (label, val) in enumerate(links):
        _add_text(slide, Inches(7.4), Inches(5.5) + Inches(i * 0.35), Inches(2), Inches(0.3),
                  label, font_size=12, color=LIGHT_GRAY)
        _add_text(slide, Inches(9.2), Inches(5.5) + Inches(i * 0.35), Inches(3), Inches(0.3),
                  val, font_size=12, color=ACCENT)

    _slide_number(slide, 11)


# ── SLIDE 12 – CALL TO ACTION ──────────────────────────────────────────────
def slide_cta():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)

    # Decorative accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(5.9), Inches(1.5), Inches(1.5), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

    _add_text(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.0),
              "Let's build transparent\nESG ratings together.",
              font_size=44, bold=True, alignment=PP_ALIGN.CENTER)

    _add_text(slide, Inches(2), Inches(3.5), Inches(9.3), Inches(0.6),
              "SustainabilitySignals is looking for early partners, advisors, and collaborators\n"
              "who believe ESG ratings should show their work.",
              font_size=17, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Contact buttons
    _add_rect(slide, Inches(3.3), Inches(4.6), Inches(6.7), Inches(1.8), CARD_BG, ACCENT, Pt(1))

    contacts = [
        ("📧  nikhilreddy.gogu@student.ams.ac.be", Inches(3.7)),
        ("🌐  nikhil.chat", Inches(4.85)),
        ("💼  linkedin.com/in/nikhilgogu", Inches(5.55)),
    ]
    for text, y in contacts:
        _add_text(slide, Inches(4.5), y, Inches(5), Inches(0.5),
                  text, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    _add_text(slide, Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
              "SustainabilitySignals  ·  © 2026  ·   Antwerpen, Belgium",
              font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    _slide_number(slide, 12)


# ── Build deck ──────────────────────────────────────────────────────────────
slide_title()
slide_problem()
slide_solution()
slide_dq_engine()
slide_ai_pipeline()
slide_coverage()
slide_features()
slide_tech()
slide_traction()
slide_roadmap()
slide_team()
slide_cta()

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "SustainabilitySignals_Pitchdeck.pptx")
out = os.path.normpath(out)
prs.save(out)
print(f"✅ Pitch deck saved to: {out}")
