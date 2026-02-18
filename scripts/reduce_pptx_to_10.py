#!/usr/bin/env python3
"""
Reduce SustainabilitySignals.pptx from 12 slides to 10 by removing
slides 8 (Technology — redundant with AI & ML slide 5) and 12 (CTA —
contact info already on Team slide 11), then convert to PDF.
"""

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import copy
import os
import subprocess
import sys

SRC = os.path.join(os.path.dirname(__file__), "..", "SustainabilitySignals.pptx")
DST_PPTX = os.path.join(os.path.dirname(__file__), "..", "SustainabilitySignals_10slides.pptx")
DST_PDF = os.path.join(os.path.dirname(__file__), "..", "SustainabilitySignals_10slides.pdf")


def delete_slide(prs, slide_index):
    """Delete a slide by index (0-based) from the presentation."""
    rId = None
    slide = prs.slides[slide_index]
    # Find the relationship ID for this slide
    for rel in prs.part.rels.values():
        if rel.target_part == slide.part:
            rId = rel.rId
            break
    if rId is None:
        raise ValueError(f"Could not find relationship for slide {slide_index}")
    
    # Remove from slide list XML
    from lxml import etree
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
             'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
    sldIdLst = prs.part._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
    for sldId in list(sldIdLst):
        if sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id') == rId:
            sldIdLst.remove(sldId)
            break
    
    # Remove the relationship
    prs.part.rels.pop(rId)


def update_slide_numbers(prs):
    """Update page number text boxes on each slide."""
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    # Find standalone number text boxes (page numbers)
                    if text.isdigit() and int(text) <= 12:
                        for run in paragraph.runs:
                            run.text = str(i)
                        break


def main():
    prs = Presentation(os.path.abspath(SRC))
    print(f"Original slide count: {len(prs.slides)}")
    
    # Delete slide 12 first (0-indexed: 11), then slide 8 (0-indexed: 7)
    # Must delete from the end to preserve indices
    slides_to_remove = [11, 7]  # 0-indexed: slide 12 (CTA) and slide 8 (Technology)
    
    for idx in sorted(slides_to_remove, reverse=True):
        slide = prs.slides[idx]
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        title = p.text.strip()[:60]
                        break
            if title:
                break
        print(f"  Removing slide {idx + 1}: {title}")
        delete_slide(prs, idx)
    
    print(f"New slide count: {len(prs.slides)}")
    
    # Update page numbers
    update_slide_numbers(prs)
    
    # Print final slide list
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip() and not p.text.strip().isdigit():
                        title = p.text.strip()[:60]
                        break
            if title:
                break
        print(f"  Slide {i}: {title}")
    
    out_pptx = os.path.abspath(DST_PPTX)
    prs.save(out_pptx)
    print(f"\nSaved 10-slide PPTX: {out_pptx}")
    
    # --- Convert to PDF ---
    out_pdf = os.path.abspath(DST_PDF)
    convert_to_pdf(out_pptx, out_pdf)


def convert_to_pdf(pptx_path, pdf_path):
    """Convert PPTX to PDF using PowerShell COM automation."""
    ps_script = f'''
$pptx = "{pptx_path.replace(os.sep, '/')}"
$pdf  = "{pdf_path.replace(os.sep, '/')}"

try {{
    $ppt = New-Object -ComObject PowerPoint.Application
    # msoTrue = -1, msoFalse = 0; ppSaveAsPDF = 32
    $presentation = $ppt.Presentations.Open($pptx, -1, 0, 0)
    $presentation.SaveAs($pdf, 32)
    $presentation.Close()
    $ppt.Quit()
    [System.Runtime.Interopservices.Marshal]::ReleaseComObject($ppt) | Out-Null
    Write-Host "PDF saved via PowerPoint: $pdf"
}} catch {{
    Write-Host "PowerPoint COM failed: $_"
    Write-Host "Trying LibreOffice..."
    
    $libreoffice = $null
    $paths = @(
        "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
        "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
        "$env:LOCALAPPDATA\\Programs\\LibreOffice\\program\\soffice.exe"
    )
    foreach ($p in $paths) {{
        if (Test-Path $p) {{ $libreoffice = $p; break }}
    }}
    
    if ($libreoffice) {{
        $outDir = Split-Path $pdf -Parent
        & $libreoffice --headless --convert-to pdf --outdir $outDir $pptx
        Write-Host "PDF saved via LibreOffice: $pdf"
    }} else {{
        Write-Host "WARN: Neither PowerPoint nor LibreOffice found for PPTX->PDF."
        Write-Host "The 10-slide PPTX has been saved. To get PDF, either:"
        Write-Host "  1) Install LibreOffice and re-run"
        Write-Host "  2) Open the PPTX and File > Save As > PDF"
    }}
}}
'''
    print("Converting PPTX to PDF...")
    result = subprocess.run(
        ["powershell", "-NoProfile", "-Command", ps_script],
        capture_output=True, text=True, timeout=120
    )
    print(result.stdout.strip())
    if result.stderr.strip():
        print(f"STDERR: {result.stderr.strip()}")


if __name__ == "__main__":
    main()
